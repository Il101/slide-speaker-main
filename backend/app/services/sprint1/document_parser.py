"""
Sprint 1: Document parsing service with concept extraction

⚠️  DEPRECATED: This module is deprecated!
    Use the new pipeline (intelligent_optimized.py) instead.
    Set USE_NEW_PIPELINE=true (it's now the default).
    
    This file is kept for backward compatibility only.
"""
import warnings
warnings.warn(
    "document_parser.py is deprecated. Use new pipeline with USE_NEW_PIPELINE=true",
    DeprecationWarning,
    stacklevel=2
)
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging
import fitz  # PyMuPDF
from PIL import Image
import cv2
import numpy as np
import easyocr
import uuid
from functools import lru_cache
import hashlib

from ...core.exceptions import ParsingError, FileProcessingError
from ...models.schemas import Slide, Manifest, SlideElement, Cue, ActionType, Timeline
from ..sprint2.concept_extractor import extract_slide_concepts, SlideConcepts

logger = logging.getLogger(__name__)

# Constants for fallback elements
ORIG_W, ORIG_H = 1600, 900

def _placeholder_element():
    """Create a placeholder element for slides with no detected content"""
    return {
        "id": "slide_area",
        "type": "placeholder", 
        "text": "",
        "bbox": [0, 0, ORIG_W, ORIG_H],
        "confidence": 1.0,
        "source": "fallback"
    }

def _ocr_enabled() -> bool:
    """Check if OCR is enabled in configuration"""
    try:
        from ...core.config import settings
        provider = settings.OCR_PROVIDER.lower()
        return provider in ["easyocr", "google", "vision", "enhanced_vision", "paddle"]
    except Exception:
        return True

def _extract_with_ocr(slide_png: Path) -> List[Dict[str, Any]]:
    """Extract elements using configured OCR provider"""
    try:
        from ...services.provider_factory import extract_elements_from_pages
        
        # Extract elements using configured OCR provider
        elements_data = extract_elements_from_pages([str(slide_png)])
        
        if elements_data and elements_data[0]:
            logger.info(f"OCR extracted {len(elements_data[0])} elements from {slide_png}")
            return elements_data[0]
        else:
            logger.warning(f"OCR returned no elements for {slide_png}")
            return []
            
    except Exception as e:
        logger.error(f"OCR extraction failed for {slide_png}: {e}")
        raise

def _extract_from_vector_layers(slide_png: Path) -> List[Dict[str, Any]]:
    """Extract elements from vector layers (placeholder for future implementation)"""
    # TODO: Implement vector layer extraction for PPTX files
    return []

def build_slide_elements(slide_png: Path) -> List[Dict[str, Any]]:
    """Build slide elements with fallback guarantees"""
    elements: List[Dict[str, Any]] = []
    
    # Try vector layer extraction first
    elements += _extract_from_vector_layers(slide_png)
    
    # If no elements and OCR is enabled, try OCR
    if not elements and _ocr_enabled():
        try:
            elements += _extract_with_ocr(slide_png)
        except Exception as e:
            logger.warning("OCR failed %s: %s", slide_png, e)
    
    # Guarantee at least one element
    if not elements:
        elements = [_placeholder_element()]
        logger.info(f"Using fallback placeholder element for {slide_png}")
    
    return elements

# OCR cache for repeated processing of same images
@lru_cache(maxsize=100)
def _get_image_hash(image_path: str) -> str:
    """Generate hash for image file to use as cache key"""
    try:
        with open(image_path, 'rb') as f:
            return hashlib.md5(f.read()).hexdigest()
    except Exception:
        return hashlib.md5(image_path.encode()).hexdigest()

@lru_cache(maxsize=50)
def _cached_ocr_detection(image_hash: str, image_path: str) -> List[SlideElement]:
    """Cached OCR detection to avoid reprocessing same images"""
    logger.info(f"Running OCR detection for cached image: {image_path}")
    
    # Load image
    image = cv2.imread(image_path)
    image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
    
    # Initialize OCR reader (this should be cached at class level)
    ocr_reader = easyocr.Reader(['en', 'ru'])
    
    # Run OCR
    results = ocr_reader.readtext(image_rgb)
    
    elements = []
    for i, (bbox, text, confidence) in enumerate(results):
        if confidence > 0.5:  # Filter low confidence results
            # Convert bbox to [x, y, width, height] format
            x_coords = [point[0] for point in bbox]
            y_coords = [point[1] for point in bbox]
            x_min, x_max = min(x_coords), max(x_coords)
            y_min, y_max = min(y_coords), max(y_coords)
            
            element = SlideElement(
                id=f"element_{uuid.uuid4().hex[:8]}",
                type="text",
                bbox=[int(x_min), int(y_min), int(x_max - x_min), int(y_max - y_min)],
                text=text.strip(),
                confidence=confidence
            )
            elements.append(element)
    
    return elements

class DocumentParser:
    """Base class for document parsing"""
    
    def __init__(self, file_path: Path):
        self.file_path = file_path
        self.output_dir = None
        
    async def parse(self) -> Manifest:
        """Parse document and return manifest"""
        raise NotImplementedError
        
    async def extract_slides(self) -> List[Slide]:
        """Extract slides from document"""
        raise NotImplementedError
        
    async def detect_elements(self, slide_image_path: Path) -> List[Dict[str, Any]]:
        """Detect text elements and their bounding boxes"""
        raise NotImplementedError

class PPTXParser(DocumentParser):
    """PowerPoint presentation parser"""
    
    def __init__(self, file_path: Path):
        super().__init__(file_path)
        self.ocr_reader = easyocr.Reader(['en', 'ru'])
        
    async def parse(self) -> Manifest:
        """Parse PPTX file and create manifest"""
        try:
            logger.info(f"Parsing PPTX file: {self.file_path}")
            
            # Create output directory
            self.output_dir = self.file_path.parent
            self.output_dir.mkdir(exist_ok=True)
            (self.output_dir / "slides").mkdir(exist_ok=True)
            (self.output_dir / "audio").mkdir(exist_ok=True)
            
            slides = await self.extract_slides()
            
            # Generate timeline with basic rules
            timeline = Timeline(
                rules=[
                    {"action_type": "highlight", "duration": 2.0, "priority": 1},
                    {"action_type": "underline", "duration": 1.5, "priority": 2},
                    {"action_type": "laser_move", "duration": 0.5, "priority": 3}
                ],
                default_duration=2.0,
                transition_duration=0.5,
                slide_change_events=[]
            )
            
            return Manifest(
                slides=slides,
                timeline=timeline,
                metadata={
                    "source_file": str(self.file_path),
                    "parser": "pptx",
                    "total_slides": len(slides)
                }
            )
            
        except Exception as e:
            logger.error(f"Error parsing PPTX: {e}")
            raise ParsingError(f"Failed to parse PPTX file: {e}")
    
    async def extract_slides(self) -> List[Slide]:
        """Extract slides from PPTX"""
        logger.info("Extracting slides from PPTX")
        
        try:
            from pptx import Presentation
        except ImportError:
            logger.error("python-pptx not installed")
            raise ParsingError("python-pptx library not available")
        
        slides = []
        prs = Presentation(self.file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            # Extract slide elements first
            elements = await self.extract_slide_elements(slide)
            
            # Convert slide to image (placeholder - would need additional library)
            # For now, we'll create a placeholder image
            slide_image_path = self.output_dir / "slides" / f"{slide_num + 1:03d}.png"
            await self.create_placeholder_slide_image(slide_image_path, elements)
            
            # Detect text elements using OCR
            detected_elements = await self.detect_elements(slide_image_path)
            
            # Generate basic cues based on elements
            cues = await self.generate_cues(detected_elements)
            
            slide_data = Slide(
                id=slide_num + 1,
                image=f"/assets/{self.file_path.parent.name}/slides/{slide_num + 1:03d}.png",
                audio=f"/assets/{self.file_path.parent.name}/audio/{slide_num + 1:03d}.mp3",
                elements=detected_elements,
                cues=cues
            )
            slides.append(slide_data)
        
        return slides
    
    async def extract_slide_elements(self, slide) -> List[Dict[str, Any]]:
        """Extract elements from PowerPoint slide"""
        elements = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                # Get shape position and size
                left = int(shape.left)
                top = int(shape.top)
                width = int(shape.width)
                height = int(shape.height)
                
                elements.append({
                    "type": "text",
                    "bbox": [left, top, width, height],
                    "text": shape.text.strip(),
                    "confidence": 1.0
                })
        
        return elements
    
    async def create_placeholder_slide_image(self, image_path: Path, elements: List[Dict[str, Any]]):
        """Create a placeholder slide image"""
        # Create a white background
        img = Image.new('RGB', (1920, 1080), 'white')
        
        # Draw elements (simplified)
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except (OSError, IOError, FileNotFoundError):
            font = ImageFont.load_default()
        
        for element in elements:
            if element["type"] == "text":
                x, y, w, h = element["bbox"]
                # Scale coordinates to image size
                x = int(x * 1920 / 914400)  # PowerPoint units to pixels
                y = int(y * 1080 / 685800)
                w = int(w * 1920 / 914400)
                h = int(h * 1080 / 685800)
                
                # Draw text
                draw.text((x, y), element["text"], fill='black', font=font)
        
        img.save(image_path)
    
    async def detect_elements(self, slide_image_path: Path) -> List[SlideElement]:
        """Detect text elements using configured OCR provider with fallback guarantees"""
        logger.info(f"Detecting elements in slide: {slide_image_path}")
        
        # Use the new build_slide_elements function with fallback guarantees
        elements_data = build_slide_elements(slide_image_path)
        
        # Convert to SlideElement objects
        elements = []
        for element_data in elements_data:
            element = SlideElement(
                id=element_data.get("id", f"element_{uuid.uuid4().hex[:8]}"),
                type=element_data.get("type", "text"),
                bbox=element_data.get("bbox", [0, 0, 100, 50]),
                text=element_data.get("text", ""),
                confidence=element_data.get("confidence", 0.9)
            )
            elements.append(element)
        
        logger.info(f"Detected {len(elements)} elements for slide: {slide_image_path}")
        return elements
    
    async def extract_slide_concepts(self, elements: List[SlideElement]) -> SlideConcepts:
        """Extract concepts from slide elements"""
        logger.info(f"Extracting concepts from {len(elements)} elements")
        
        # Convert SlideElement objects to dict format for concept extraction
        elements_dict = []
        for element in elements:
            elements_dict.append({
                "id": element.id,
                "type": element.type,
                "bbox": element.bbox,
                "text": element.text,
                "confidence": element.confidence
            })
        
        # Extract concepts using the concept extractor
        concepts = extract_slide_concepts(elements_dict)
        
        logger.info(f"Extracted concepts: title='{concepts.title}', {len(concepts.key_theses)} theses, {len(concepts.terms_to_define)} terms")
        return concepts
    
    async def _detect_elements_easyocr(self, slide_image_path: Path) -> List[SlideElement]:
        """Detect text elements using EasyOCR (fallback) with caching"""
        logger.info(f"Detecting elements using EasyOCR: {slide_image_path}")
        
        # Check cache first
        image_hash = _get_image_hash(str(slide_image_path))
        try:
            cached_elements = _cached_ocr_detection(image_hash, str(slide_image_path))
            logger.info(f"Using cached OCR results for {slide_image_path}")
            return cached_elements
        except Exception as e:
            logger.warning(f"Cache miss or error, running fresh OCR: {e}")
        
        # Fallback to direct OCR if cache fails
        image = cv2.imread(str(slide_image_path))
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        
        # Run OCR
        results = self.ocr_reader.readtext(image_rgb)
        
        elements = []
        for i, (bbox, text, confidence) in enumerate(results):
            if confidence > 0.5:  # Filter low confidence results
                # Convert bbox to [x, y, width, height] format
                x_coords = [point[0] for point in bbox]
                y_coords = [point[1] for point in bbox]
                x_min, x_max = min(x_coords), max(x_coords)
                y_min, y_max = min(y_coords), max(y_coords)
                
                element = SlideElement(
                    id=f"element_{uuid.uuid4().hex[:8]}",
                    type="text",
                    bbox=[int(x_min), int(y_min), int(x_max - x_min), int(y_max - y_min)],
                    text=text.strip(),
                    confidence=confidence
                )
                elements.append(element)
        
        return elements
    
    async def generate_cues(self, elements: List[SlideElement]) -> List[Cue]:
        """Generate basic visual cues based on detected elements"""
        cues = []
        current_time = 0.5  # Start after 0.5 seconds
        
        for i, element in enumerate(elements):
            if element.type == "text" and element.text:
                # Highlight effect
                highlight_cue = Cue(
                    cue_id=f"cue_{uuid.uuid4().hex[:8]}",
                    t0=current_time,
                    t1=current_time + 2.0,
                    action=ActionType.HIGHLIGHT,
                    bbox=element.bbox,
                    element_id=element.id
                )
                cues.append(highlight_cue)
                
                # Underline effect
                underline_cue = Cue(
                    cue_id=f"cue_{uuid.uuid4().hex[:8]}",
                    t0=current_time + 0.5,
                    t1=current_time + 2.5,
                    action=ActionType.UNDERLINE,
                    bbox=[element.bbox[0], element.bbox[1] + element.bbox[3] - 4, element.bbox[2], 4],
                    element_id=element.id
                )
                cues.append(underline_cue)
                
                # Laser move to next element
                if i < len(elements) - 1:
                    next_element = elements[i + 1]
                    laser_cue = Cue(
                        cue_id=f"cue_{uuid.uuid4().hex[:8]}",
                        t0=current_time + 2.0,
                        t1=current_time + 2.5,
                        action=ActionType.LASER_MOVE,
                        to=[next_element.bbox[0] + next_element.bbox[2] // 2, 
                            next_element.bbox[1] + next_element.bbox[3] // 2]
                    )
                    cues.append(laser_cue)
                
                current_time += 3.0  # Move to next element after 3 seconds
        
        return cues

class PDFParser(DocumentParser):
    """PDF document parser"""
    
    def __init__(self, file_path: Path):
        super().__init__(file_path)
        self.ocr_reader = easyocr.Reader(['en', 'ru'])
        
    async def parse(self) -> Manifest:
        """Parse PDF file and create manifest"""
        try:
            logger.info(f"Parsing PDF file: {self.file_path}")
            
            # Create output directory
            self.output_dir = self.file_path.parent
            self.output_dir.mkdir(exist_ok=True)
            (self.output_dir / "slides").mkdir(exist_ok=True)
            (self.output_dir / "audio").mkdir(exist_ok=True)
            
            slides = await self.extract_slides()
            
            # Generate timeline with basic rules
            timeline = Timeline(
                rules=[
                    {"action_type": "highlight", "duration": 2.0, "priority": 1},
                    {"action_type": "underline", "duration": 1.5, "priority": 2},
                    {"action_type": "laser_move", "duration": 0.5, "priority": 3}
                ],
                default_duration=2.0,
                transition_duration=0.5,
                slide_change_events=[]
            )
            
            return Manifest(
                slides=slides,
                timeline=timeline,
                metadata={
                    "source_file": str(self.file_path),
                    "parser": "pdf",
                    "total_slides": len(slides)
                }
            )
            
        except Exception as e:
            logger.error(f"Error parsing PDF: {e}")
            raise ParsingError(f"Failed to parse PDF file: {e}")
    
    async def extract_slides(self) -> List[Slide]:
        """Extract pages from PDF as slides"""
        logger.info("Extracting pages from PDF")
        
        slides = []
        doc = fitz.open(self.file_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Convert page to image
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Save slide image
            slide_image_path = self.output_dir / "slides" / f"{page_num + 1:03d}.png"
            with open(slide_image_path, "wb") as f:
                f.write(img_data)
            
            # Detect elements
            elements = await self.detect_elements(slide_image_path)
            
            # Generate basic cues based on elements
            cues = await self.generate_cues(elements)
            
            slide = Slide(
                id=page_num + 1,
                image=f"/assets/{self.file_path.parent.name}/slides/{page_num + 1:03d}.png",
                audio=f"/assets/{self.file_path.parent.name}/audio/{page_num + 1:03d}.mp3",
                elements=elements,
                cues=cues
            )
            slides.append(slide)
        
        doc.close()
        return slides
    
    async def detect_elements(self, slide_image_path: Path) -> List[SlideElement]:
        """Detect text elements using configured OCR provider with fallback guarantees"""
        logger.info(f"Detecting elements in page: {slide_image_path}")
        
        # Use the new build_slide_elements function with fallback guarantees
        elements_data = build_slide_elements(slide_image_path)
        
        # Convert to SlideElement objects
        elements = []
        for element_data in elements_data:
            element = SlideElement(
                id=element_data.get("id", f"element_{uuid.uuid4().hex[:8]}"),
                type=element_data.get("type", "text"),
                bbox=element_data.get("bbox", [0, 0, 100, 50]),
                text=element_data.get("text", ""),
                confidence=element_data.get("confidence", 0.9)
            )
            elements.append(element)
        
        logger.info(f"Detected {len(elements)} elements for page: {slide_image_path}")
        return elements
    
    async def extract_slide_concepts(self, elements: List[SlideElement]) -> SlideConcepts:
        """Extract concepts from slide elements"""
        logger.info(f"Extracting concepts from {len(elements)} elements")
        
        # Convert SlideElement objects to dict format for concept extraction
        elements_dict = []
        for element in elements:
            elements_dict.append({
                "id": element.id,
                "type": element.type,
                "bbox": element.bbox,
                "text": element.text,
                "confidence": element.confidence
            })
        
        # Extract concepts using the concept extractor
        concepts = extract_slide_concepts(elements_dict)
        
        logger.info(f"Extracted concepts: title='{concepts.title}', {len(concepts.key_theses)} theses, {len(concepts.terms_to_define)} terms")
        return concepts
    
    async def _detect_elements_easyocr(self, slide_image_path: Path) -> List[SlideElement]:
        """Detect text elements using EasyOCR (fallback) with caching"""
        logger.info(f"Detecting elements using EasyOCR: {slide_image_path}")
        
        # Check cache first
        image_hash = _get_image_hash(str(slide_image_path))
        try:
            cached_elements = _cached_ocr_detection(image_hash, str(slide_image_path))
            logger.info(f"Using cached OCR results for {slide_image_path}")
            return cached_elements
        except Exception as e:
            logger.warning(f"Cache miss or error, running fresh OCR: {e}")
        
        # Fallback to direct OCR if cache fails
        image = cv2.imread(str(slide_image_path))
        image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
        
        # Run OCR
        results = self.ocr_reader.readtext(image_rgb)
        
        elements = []
        for i, (bbox, text, confidence) in enumerate(results):
            if confidence > 0.5:  # Filter low confidence results
                # Convert bbox to [x, y, width, height] format
                x_coords = [point[0] for point in bbox]
                y_coords = [point[1] for point in bbox]
                x_min, x_max = min(x_coords), max(x_coords)
                y_min, y_max = min(y_coords), max(y_coords)
                
                element = SlideElement(
                    id=f"element_{uuid.uuid4().hex[:8]}",
                    type="text",
                    bbox=[int(x_min), int(y_min), int(x_max - x_min), int(y_max - y_min)],
                    text=text.strip(),
                    confidence=confidence
                )
                elements.append(element)
        
        return elements
    
    async def generate_cues(self, elements: List[SlideElement]) -> List[Cue]:
        """Generate basic visual cues based on detected elements"""
        cues = []
        current_time = 0.5  # Start after 0.5 seconds
        
        for i, element in enumerate(elements):
            if element.type == "text" and element.text:
                # Highlight effect
                highlight_cue = Cue(
                    cue_id=f"cue_{uuid.uuid4().hex[:8]}",
                    t0=current_time,
                    t1=current_time + 2.0,
                    action=ActionType.HIGHLIGHT,
                    bbox=element.bbox,
                    element_id=element.id
                )
                cues.append(highlight_cue)
                
                # Underline effect
                underline_cue = Cue(
                    cue_id=f"cue_{uuid.uuid4().hex[:8]}",
                    t0=current_time + 0.5,
                    t1=current_time + 2.5,
                    action=ActionType.UNDERLINE,
                    bbox=[element.bbox[0], element.bbox[1] + element.bbox[3] - 4, element.bbox[2], 4],
                    element_id=element.id
                )
                cues.append(underline_cue)
                
                # Laser move to next element
                if i < len(elements) - 1:
                    next_element = elements[i + 1]
                    laser_cue = Cue(
                        cue_id=f"cue_{uuid.uuid4().hex[:8]}",
                        t0=current_time + 2.0,
                        t1=current_time + 2.5,
                        action=ActionType.LASER_MOVE,
                        to=[next_element.bbox[0] + next_element.bbox[2] // 2, 
                            next_element.bbox[1] + next_element.bbox[3] // 2]
                    )
                    cues.append(laser_cue)
                
                current_time += 3.0  # Move to next element after 3 seconds
        
        return cues

class ParserFactory:
    """Factory for creating document parsers"""
    
    @staticmethod
    def create_parser(file_path: Path) -> DocumentParser:
        """Create appropriate parser based on file extension"""
        extension = file_path.suffix.lower()
        
        if extension == '.pptx':
            return PPTXParser(file_path)
        elif extension == '.pdf':
            return PDFParser(file_path)
        else:
            raise FileProcessingError(f"Unsupported file format: {extension}")