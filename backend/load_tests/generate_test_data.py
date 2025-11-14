"""
Test data generator for load testing
Creates realistic test data including users, presentations, and lessons
"""

import asyncio
import random
import string
from pathlib import Path
from typing import List, Dict
import httpx
import io
from pptx import Presentation
from pptx.util import Inches, Pt


class TestDataGenerator:
    """Generate realistic test data for load testing"""
    
    def __init__(self, base_url: str):
        self.base_url = base_url
        self.users: List[Dict] = []
        
    def generate_random_email(self) -> str:
        """Generate random email for test user"""
        random_part = ''.join(random.choices(string.ascii_lowercase + string.digits, k=8))
        return f"loadtest_{random_part}@example.com"
    
    def generate_random_password(self) -> str:
        """Generate secure random password"""
        return ''.join(random.choices(string.ascii_letters + string.digits + "!@#$%", k=16))
    
    async def create_test_users(self, count: int = 100) -> List[Dict]:
        """Create test users via API"""
        print(f"Creating {count} test users...")
        
        async with httpx.AsyncClient(base_url=self.base_url, timeout=30.0) as client:
            for i in range(count):
                email = self.generate_random_email()
                password = self.generate_random_password()
                
                try:
                    response = await client.post("/auth/register", json={
                        "email": email,
                        "password": password,
                        "name": f"Load Test User {i+1}"
                    })
                    
                    if response.status_code == 200:
                        data = response.json()
                        self.users.append({
                            "email": email,
                            "password": password,
                            "token": data.get("access_token"),
                            "user_id": data.get("user_id")
                        })
                        
                        if (i + 1) % 10 == 0:
                            print(f"  Created {i+1}/{count} users...")
                            
                except Exception as e:
                    print(f"  Error creating user {i+1}: {e}")
        
        print(f"✓ Created {len(self.users)} test users")
        return self.users
    
    def create_test_presentation(self, num_slides: int = 10, size: str = "small") -> bytes:
        """Create a test PowerPoint presentation"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define content complexity based on size
        paragraphs_per_slide = {
            "small": 2,
            "medium": 5,
            "large": 10
        }
        
        words_per_paragraph = {
            "small": 10,
            "medium": 30,
            "large": 50
        }
        
        num_paragraphs = paragraphs_per_slide.get(size, 2)
        num_words = words_per_paragraph.get(size, 10)
        
        for i in range(num_slides):
            # Add slide
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title = slide.shapes.title
            title.text = f"Test Slide {i + 1}"
            
            # Add content
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                tf = content.text_frame
                
                for p in range(num_paragraphs):
                    if p > 0:
                        tf.add_paragraph()
                    
                    paragraph = tf.paragraphs[p]
                    # Generate Lorem Ipsum-style text
                    words = [random.choice([
                        "lorem", "ipsum", "dolor", "sit", "amet", "consectetur",
                        "adipiscing", "elit", "sed", "do", "eiusmod", "tempor",
                        "incididunt", "ut", "labore", "et", "dolore", "magna"
                    ]) for _ in range(num_words)]
                    paragraph.text = " ".join(words).capitalize() + "."
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()
    
    async def upload_test_presentations(self, presentations_per_user: int = 3) -> None:
        """Upload test presentations for all users"""
        print(f"\nUploading {presentations_per_user} presentations per user...")
        
        total_uploads = 0
        
        async with httpx.AsyncClient(base_url=self.base_url, timeout=60.0) as client:
            for i, user in enumerate(self.users):
                headers = {"Authorization": f"Bearer {user['token']}"}
                
                for j in range(presentations_per_user):
                    # Vary presentation sizes
                    size = random.choice(["small", "medium", "large"])
                    num_slides = random.randint(5, 20)
                    
                    # Create presentation
                    pptx_bytes = self.create_test_presentation(
                        num_slides=num_slides,
                        size=size
                    )
                    
                    # Upload
                    files = {
                        'file': (
                            f'test_{i}_{j}_{size}.pptx',
                            pptx_bytes,
                            'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                        )
                    }
                    
                    try:
                        response = await client.post(
                            "/upload",
                            files=files,
                            headers=headers
                        )
                        
                        if response.status_code == 200:
                            total_uploads += 1
                            
                    except Exception as e:
                        print(f"  Error uploading for user {i+1}: {e}")
                
                if (i + 1) % 10 == 0:
                    print(f"  Processed {i+1}/{len(self.users)} users ({total_uploads} uploads)...")
        
        print(f"✓ Uploaded {total_uploads} test presentations")
    
    def save_credentials(self, filepath: str = "test_users.txt"):
        """Save user credentials to file"""
        output_path = Path(filepath)
        
        with open(output_path, 'w') as f:
            f.write("# Test User Credentials\n")
            f.write("# Generated for load testing\n\n")
            
            for user in self.users:
                f.write(f"Email: {user['email']}\n")
                f.write(f"Password: {user['password']}\n")
                f.write(f"Token: {user['token']}\n")
                f.write(f"User ID: {user['user_id']}\n")
                f.write("-" * 80 + "\n")
        
        print(f"\n✓ Saved credentials to {output_path}")


async def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(description="Generate test data for load testing")
    parser.add_argument("--url", default="http://localhost:8000", help="Base URL of the API")
    parser.add_argument("--users", type=int, default=50, help="Number of test users to create")
    parser.add_argument("--presentations", type=int, default=2, help="Presentations per user")
    parser.add_argument("--skip-upload", action="store_true", help="Skip uploading presentations")
    
    args = parser.parse_args()
    
    print("=" * 80)
    print("Test Data Generator for Slide Speaker Load Testing")
    print("=" * 80)
    print(f"Target URL: {args.url}")
    print(f"Users to create: {args.users}")
    print(f"Presentations per user: {args.presentations}")
    print()
    
    generator = TestDataGenerator(args.url)
    
    # Create users
    await generator.create_test_users(count=args.users)
    
    # Upload presentations
    if not args.skip_upload:
        await generator.upload_test_presentations(presentations_per_user=args.presentations)
    
    # Save credentials
    generator.save_credentials("test_users.txt")
    
    print("\n" + "=" * 80)
    print("✓ Test data generation complete!")
    print("=" * 80)


if __name__ == "__main__":
    asyncio.run(main())
